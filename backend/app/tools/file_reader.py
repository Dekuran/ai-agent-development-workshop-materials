from pathlib import Path
import json
import csv

# Optional dependency imports
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None


REPO_ROOT = Path(__file__).resolve().parents[2]
UPLOAD_DIR = REPO_ROOT / "uploaded_files"


class FileReader:
    """
    Reads files from the uploaded_files directory in various formats.
    Supported formats: txt, md, log, csv, json, pdf, docx, pptx, png, jpg, jpeg, gif, bmp.
    Returns text content for text-based files, and a message for images/unsupported types.
    """

    def read(self, relative_path: str) -> str:
        """
        Read the file at the given relative path and return its contents as a string.
        Handles txt, md, log, csv, json, pdf, docx, pptx, and image files.
        """
        p = UPLOAD_DIR / relative_path

        print(f"[FileReader] read called with relative_path: {relative_path}")

        # SECURITY CHECK: prevent directory traversal
        if UPLOAD_DIR not in p.parents and p != UPLOAD_DIR:
            raise ValueError("Access denied")

        try:
            ext = p.suffix.lower()

            # ---- TEXT FILES ----
            if ext in [".txt", ".md", ".log", ""]:
                return p.read_text(encoding="utf-8")

            # ---- CSV ----
            elif ext == ".csv":
                with p.open("r", encoding="utf-8") as f:
                    reader = csv.reader(f)
                    return "\n".join([", ".join(row) for row in reader])

            # ---- JSON ----
            elif ext == ".json":
                with p.open("r", encoding="utf-8") as f:
                    data = json.load(f)
                return str(data)

            # ---- PDF ----
            elif ext == ".pdf":
                if PyPDF2 is None:
                    return "[FileReader] PyPDF2 not installed. Cannot read PDF files."
                with p.open("rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text = "\n".join(page.extract_text() or "" for page in reader.pages)
                return text or "[FileReader] No extractable text in PDF."

            # ---- DOCX ----
            elif ext == ".docx":
                if docx is None:
                    return "[FileReader] python-docx not installed. Cannot read DOCX files."
                doc = docx.Document(str(p))
                return "\n".join([para.text for para in doc.paragraphs])

            # ---- PPTX ----
            elif ext == ".pptx":
                if pptx is None:
                    return "[FileReader] python-pptx not installed. Cannot read PPTX files."
                pres = pptx.Presentation(str(p))
                slides = []
                for slide in pres.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slides.append(shape.text)
                return "\n".join(slides)

            # ---- IMAGES ----
            elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]:
                return f"[FileReader] {relative_path} is an image file."

            # ---- UNSUPPORTED ----
            else:
                return f"[FileReader] Unsupported file type: {ext}"

        except Exception as e:
            return f"[FileReader] Error reading file: {e}"
